"""
生成股票量化交易系统架构汇报 PPT
输出: 量化交易系统架构汇报.pptx
"""
import os, sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 颜色方案
PRIMARY = RGBColor(0x1A, 0x23, 0x3E)      # 深蓝黑
SECONDARY = RGBColor(0x21, 0x96, 0xF3)    # 蓝色
ACCENT = RGBColor(0x4C, 0xAF, 0x50)       # 绿色
WARNING = RGBColor(0xFF, 0x98, 0x00)      # 橙色
DANGER = RGBColor(0xF4, 0x43, 0x36)       # 红色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x75, 0x75, 0x75)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs)

    # Slide 2: 项目概述
    add_overview_slide(prs)

    # Slide 3: 系统架构总览
    add_architecture_slide(prs)

    # Slide 4: 数据采集模块
    add_data_slide(prs)

    # Slide 5: 因子指标库
    add_factors_slide(prs)

    # Slide 6: 信号生成引擎
    add_signals_slide(prs)

    # Slide 7: 仓位管理与风控
    add_portfolio_slide(prs)

    # Slide 8: 回测引擎
    add_backtest_slide(prs)

    # Slide 9: AI Agent
    add_agent_slide(prs)

    # Slide 10: 可视化与优化
    add_viz_slide(prs)

    # Slide 11: 项目结构
    add_structure_slide(prs)

    # Slide 12: 技术栈
    add_tech_slide(prs)

    # Slide 13: 策略缺陷与过拟合
    add_risk_slide(prs)

    # Slide 14: 总结
    add_summary_slide(prs)

    # 保存
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output = os.path.join(output_dir, "量化交易系统架构汇报.pptx")
    prs.save(output)
    print(f"PPT saved to: {output}")
    return output


def add_bg(slide, color=PRIMARY):
    """设置幻灯片背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(tf, items, font_size=14, color=WHITE, indent_level=0):
    """向已有 text_frame 添加项目符号列表"""
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.level = indent_level
    return tf


def add_card(slide, left, top, width, height, title, items, title_color=SECONDARY, bg_color=None):
    """添加卡片组件"""
    from pptx.util import Inches, Pt
    # 背景矩形
    if bg_color:
        shape = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()

    # 标题
    tf = add_textbox(slide, left + 0.15, top + 0.1, width - 0.3, 0.4,
                     title, font_size=14, color=title_color, bold=True)
    # 内容
    content_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.5), Inches(width - 0.3), Inches(height - 0.6)
    )
    ctf = content_box.text_frame
    ctf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = ctf.paragraphs[0]
            first = False
        else:
            p = ctf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY if bg_color else DARK_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
    return ctf


# ============================================================
# Slide 1: Title
# ============================================================
def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, PRIMARY)

    # 装饰线
    shape = slide.shapes.add_shape(1, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY
    shape.line.fill.background()

    add_textbox(slide, 1.5, 3.0, 10, 1.0,
                "股票量化交易系统", font_size=44, bold=True, color=WHITE)
    add_textbox(slide, 1.5, 3.9, 10, 0.6,
                "Quantitative Trading System — Architecture & Design", font_size=20, color=SECONDARY)
    add_textbox(slide, 1.5, 4.7, 10, 0.5,
                "模块化架构 | AI Agent 智能分析 | 对话式交互 | 完整回测链路", font_size=16, color=MID_GRAY)

    # 底部信息
    add_textbox(slide, 1.5, 6.2, 10, 0.4,
                "风险声明：本系统仅供教育及研究目的，不构成任何投资建议", font_size=12, color=DANGER)

    add_textbox(slide, 1.5, 6.6, 10, 0.4,
                "Stock Quant System · 2026", font_size=12, color=MID_GRAY)


# ============================================================
# Slide 2: 项目概述
# ============================================================
def add_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)

    add_slide_title(slide, "项目概述", "Project Overview")

    # 3列卡片
    cards = [
        ("🎯 目标", [
            "构建完整的A股量化回测系统",
            "支持自定义策略开发与验证",
            "AI Agent 辅助分析与决策",
            "教学用途，非实盘交易工具",
        ]),
        ("🏗️ 架构", [
            "10大功能模块，60+ Python文件",
            "模块化设计，各模块可独立调用",
            "统一数据Schema，多数据源支持",
            "Streamlit Web UI + CLI 双模式",
        ]),
        ("🤖 AI 能力", [
            "市场状态自动诊断 (ADX/RSI/波动率)",
            "策略智能推荐 (5种策略自动匹配)",
            "参数自动优化 (网格搜索+交叉验证)",
            "对话式交互 (自然语言查询分析)",
        ]),
    ]

    for i, (title, items) in enumerate(cards):
        left = 0.8 + i * 4.1
        add_card(slide, left, 1.8, 3.8, 4.5, title, items, title_color=SECONDARY,
                 bg_color=RGBColor(0x25, 0x30, 0x50))

    add_footer(slide, "项目定位：教学演示 · 不构成投资建议")


# ============================================================
# Slide 3: 系统架构总览
# ============================================================
def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "系统架构总览", "System Architecture — Data Flow")

    # 数据流图示
    flow = [
        ("📡 数据采集\nAKShare/YFinance", SECONDARY),
        ("🧹 数据清洗\n缺失值/异常值", SECONDARY),
        ("📊 因子计算\n12+指标", ACCENT),
        ("🎯 信号生成\n趋势/反转/多因子", ACCENT),
        ("💰 仓位风控\n3模型+5规则", WARNING),
        ("📈 回测引擎\n向量化模拟", WARNING),
        ("📉 可视化\n6种图表", DANGER),
        ("🔍 优化\n网格搜索", DANGER),
    ]

    for i, (label, color) in enumerate(flow):
        left = 0.3 + i * 1.6
        top = 2.2
        # 盒子
        shape = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(1.45), Inches(1.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER

        # 箭头
        if i < len(flow) - 1:
            arrow = slide.shapes.add_shape(
                2, Inches(left + 1.45), Inches(top + 0.5), Inches(0.15), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MID_GRAY
            arrow.line.fill.background()

    # 下方说明卡片
    modules_info = [
        ("数据层", "data/ — AKShare数据源 + 统一Schema + 清洗 + Parquet缓存", SECONDARY),
        ("引擎层", "factors/ + signals/ + portfolio/ — 因子→信号→仓位→风控", ACCENT),
        ("执行层", "backtest/ — 向量化回测 + 指标计算 + 滑点手续费模型", WARNING),
        ("分析层", "visualization/ + optimization/ — 图表 + 网格搜索 + 过拟合检测", DANGER),
        ("智能层", "agent/ — 市场诊断 + 策略推荐 + 对话式交互", RGBColor(0x9C, 0x27, 0xB0)),
        ("界面层", "app.py — Streamlit Web UI (3个标签页)", RGBColor(0x00, 0xBC, 0xD4)),
    ]

    for i, (name, desc, color) in enumerate(modules_info):
        top = 4.2 + i * 0.55
        # 色标
        dot = slide.shapes.add_shape(1, Inches(0.8), Inches(top + 0.05), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        add_textbox(slide, 1.1, top, 1.5, 0.4, name, font_size=13, color=color, bold=True)
        add_textbox(slide, 2.5, top, 10, 0.4, desc, font_size=12, color=LIGHT_GRAY)

    add_footer(slide, "数据流：Data → Clean → Factors → Signals → Portfolio → Backtest → Result")


# ============================================================
# Slide 4-12: Helper and remaining slides
# ============================================================

def add_slide_title(slide, title, subtitle):
    add_textbox(slide, 0.8, 0.4, 11, 0.6, title, font_size=32, bold=True, color=WHITE)
    add_textbox(slide, 0.8, 1.0, 11, 0.4, subtitle, font_size=14, color=MID_GRAY)
    # 分割线
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.45), Inches(5), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY
    shape.line.fill.background()


def add_footer(slide, text):
    add_textbox(slide, 0.8, 6.9, 11, 0.4, "⚠️ " + text, font_size=10, color=MID_GRAY)


# Slide 4: Data Layer
def add_data_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "数据采集模块", "Data Layer — AKShare + Unified Schema + Clean + Cache")

    # 左侧：架构说明
    items = [
        "AKShareSource: 免费A股日线下载（前复权/后复权/不复权）",
        "YFinanceSource: 美股数据（备用）",
        "column_mapper: 中文列名 → 英文标准列名",
        "统一Schema: [symbol, date, open, high, low, close, volume, amount]",
        "",
        "DataCleaner 清洗管线:",
        "  → forward_fill_ohlcv() 前向填充（max_gap=3天）",
        "  → remove_outliers_by_zscore() Z-score异常剔除",
        "  → remove_price_anomalies() 零价/负价/极端价过滤",
        "",
        "CacheManager: 内存 + Parquet双层缓存，TTL过期策略",
    ]
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    for i, item in enumerate(items):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = item
        p.font.size = Pt(13 if item else 10)
        p.font.color.rgb = WHITE if item else RGBColor(0x44, 0x44, 0x44)
        p.font.name = "Consolas" if item.startswith("  ") else "Microsoft YaHei"
        p.space_after = Pt(2)

    # 右侧：Schema 表格
    schema_data = [
        ("列名", "类型", "AKShare源", "说明"),
        ("symbol", "str", "股票代码", "标的代码"),
        ("date", "datetime", "日期", "交易日"),
        ("open", "float", "开盘", "开盘价"),
        ("high", "float", "最高", "最高价"),
        ("low", "float", "最低", "最低价"),
        ("close", "float", "收盘", "收盘价"),
        ("volume", "int", "成交量", "成交量(股)"),
        ("amount", "float", "成交额", "成交额(元)"),
    ]

    table_top = 1.8
    tbl = slide.shapes.add_table(len(schema_data), 4, Inches(7.3), Inches(table_top),
                                  Inches(5.2), Inches(3.5)).table
    for r, row in enumerate(schema_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Microsoft YaHei"
                p.font.bold = (r == 0)
                p.font.color.rgb = PRIMARY if r == 0 else WHITE
                p.alignment = PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SECONDARY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x25, 0x30, 0x50)

    add_footer(slide, "数据源：AKShare免费接口 · 本地Parquet缓存")


# Slide 5: Factors
def add_factors_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "因子指标库", "Factor Library — 12+ Technical Indicators + Registry Pattern")

    factors = [
        ("技术指标 technical.py", SECONDARY, [
            "SMAFactor(window) — 简单移动均线",
            "EMAFactor(window) — 指数移动均线",
            "MACDFactor(fast,slow,signal) — MACD线/信号线/柱",
            "RSIFactor(window) — 相对强弱指数(0-100)",
            "BollingerFactor(window,std) — 上轨/中轨/下轨/%b",
            "ATRFactor(window) — 平均真实波幅",
        ]),
        ("动量因子 momentum.py", ACCENT, [
            "RateOfChangeFactor(window) — ROC变化率",
            "MomentumFactor(window) — 价格动量",
            "RelativeStrengthFactor(window) — 相对强弱",
        ]),
        ("波动率 volatility.py", WARNING, [
            "HistoricalVolatilityFactor — 历史波动率(年化)",
            "DailyReturnFactor — 日收益率",
            "VolumeWeightedPriceFactor — VWAP偏离度",
        ]),
    ]

    for i, (title, color, items) in enumerate(factors):
        top = 1.8 + i * 1.75
        add_card(slide, 0.8, top, 4.5, 1.55, title, items, title_color=color,
                 bg_color=RGBColor(0x25, 0x30, 0x50))

    # 右侧：注册中心
    add_card(slide, 5.8, 1.8, 7.0, 1.55,
             "FactorRegistry 注册中心", [
                 "装饰器注册: @FactorRegistry.register",
                 "手动注册: FactorRegistry.add(name, cls)",
                 "创建实例: FactorRegistry.create(name, params)",
                 "列出因子: FactorRegistry.list_factors() → 12个已注册",
                 "支持用户自定义因子热插拔，无需修改核心代码",
             ], title_color=RGBColor(0x9C, 0x27, 0xB0), bg_color=RGBColor(0x25, 0x30, 0x50))

    # 代码示例
    code_box = slide.shapes.add_textbox(Inches(5.8), Inches(3.7), Inches(7.0), Inches(2.5))
    ctf = code_box.text_frame
    ctf.word_wrap = True
    code_lines = [
        "@FactorRegistry.register",
        "class MyMomentumFactor(Factor):",
        '    def compute(self, data):',
        "        return data['close'].pct_change(20)",
        "",
        "# 使用",
        "factor = FactorRegistry.create('SMAFactor', {'window': 20})",
        "result = factor.compute(ohlcv_data)",
    ]
    for i, line in enumerate(code_lines):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xA0, 0xD0, 0xFF)
        p.font.name = "Consolas"
        p.space_after = Pt(1)

    add_footer(slide, "因子基于历史价格计算 · 不具备预测能力 · 注册中心支持无限扩展")


# Slide 6: Signals
def add_signals_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "信号生成引擎", "Signal Engine — Trend + Reversal + Multi-Factor + Filter Chain")

    cards = [
        ("趋势跟踪", SECONDARY, [
            "DualMACrossoverSignal",
            "  快线(5)上穿慢线(20) → 买入",
            "  快线下穿慢线 → 卖出",
            "MACDSignalGenerator",
            "  MACD上穿信号线 → 金叉买入",
            "  下穿 → 死叉卖出",
        ]),
        ("均值回归", ACCENT, [
            "RSIReversalSignal",
            "  RSI<30且回升 → 超卖买入",
            "  RSI>70且回落 → 超买卖出",
            "BollingerReversalSignal",
            "  价格触下轨回升 → 买入",
            "  价格触上轨回落 → 卖出",
        ]),
        ("多因子打分+过滤", WARNING, [
            "MultiFactorScorer",
            "  Z-score归一化 + 加权求和",
            "  score > buy_threshold → 买入",
            "  score < sell_threshold → 卖出",
            "FilterChain (AND/OR/NOT)",
            "  min_volume: 最低成交量",
            "  volatility_filter: 波动率上限",
            "  min_bars_between_trades: 最小间隔",
        ]),
    ]

    for i, (title, color, items) in enumerate(cards):
        add_card(slide, 0.8 + i * 4.1, 1.8, 3.8, 3.8, title, items, title_color=color,
                 bg_color=RGBColor(0x25, 0x30, 0x50))

    add_footer(slide, "信号约定: +1买入 / -1卖出 / 0持仓不变 / NaN空仓 · 保留信号强度连续值")


# Slide 7: Portfolio & Risk
def add_portfolio_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "仓位管理与风控", "Portfolio & Risk — 3 Position Models + 5 Risk Rules")

    # 左侧：仓位模型
    sizing = [
        ("固定仓位 FixedFractionSizer", SECONDARY, [
            "fraction=0.10: 每次用10%资金",
            "max_fraction=0.25: 上限25%",
            "目标股数 = 资金×比例/价格",
            "适合: 简单直接，新手友好",
        ]),
        ("ATR动态 ATRPositionSizer", ACCENT, [
            "risk_pct=0.02: 单笔风险2%",
            "atr_window=14, multiplier=2",
            "目标股数 = 资金×风险/(ATR×倍数)",
            "适合: 自动适配波动率环境",
        ]),
        ("凯利公式 KellyPositionSizer", WARNING, [
            "f* = p − (1−p)/b",
            "fraction=0.25: 1/4凯利(保守)",
            "胜率×盈亏比 → 最优仓位",
            "适合: 有历史统计基础的策略",
        ]),
    ]
    for i, (title, color, items) in enumerate(sizing):
        add_card(slide, 0.8, 1.8 + i * 1.7, 6.0, 1.55, title, items, title_color=color,
                 bg_color=RGBColor(0x25, 0x30, 0x50))

    # 右侧：风控规则
    risk_rules = [
        "1. FixedStopLoss(7%): 单笔亏损>7% → 强制平仓",
        "2. TrailingStop(15%): 从最高点回撤>15% → 平仓",
        "3. DailyLossLimit(3%): 当日亏损>3% → 停止交易",
        "4. MaxDrawdownLimit(25%): 总回撤>25% → 全部平仓",
        "5. ConsecutiveLossLimit(5): 连续亏损5次 → 暂停",
    ]

    add_card(slide, 7.2, 1.8, 5.5, 1.55,
             "RiskManager 风控管理器 (组合模式)", risk_rules,
             title_color=DANGER, bg_color=RGBColor(0x25, 0x30, 0x50))

    # 风控流程
    flow_box = slide.shapes.add_textbox(Inches(7.2), Inches(3.7), Inches(5.5), Inches(2.5))
    ctf = flow_box.text_frame
    ctf.word_wrap = True
    flow_lines = [
        "风控流程:",
        "",
        "  开仓请求 → 风控链检查 →",
        "    ↓                    ↓",
        "  [止损检查]        [日内限制]",
        "    ↓                    ↓",
        "  [固定7%]          [日亏3%]",
        "    ↓                    ↓",
        "  [移动15%]         [总回撤25%]",
        "    ↓                    ↓",
        "  任一拒绝 → 拦截交易",
    ]
    for i, line in enumerate(flow_lines):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xFF, 0xCC, 0x80) if "→" in line else LIGHT_GRAY
        p.font.name = "Consolas"
        p.space_after = Pt(1)

    add_footer(slide, "风控不能完全消除风险 · 极端行情(一字跌停)下止损可能无法执行")


# Slide 8: Backtest
def add_backtest_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "回测引擎", "Backtest Engine — Vectorized Simulation + Slippage + Metrics")

    # 回测流程
    bt_items = [
        "SimpleBacktestEngine — 纯numpy向量化实现，遵守A股T+1规则",
        "",
        "逐日模拟流程:",
        "  ① 继承前日持仓/现金",
        "  ② 风控检查：止损→移动止损→日内限制→回撤限制",
        "  ③ 信号处理：买入开仓 / 卖出平仓",
        "  ④ 滑点+手续费计算（A股整手100股）",
        "  ⑤ 更新权益 = 现金 + 持仓×当日收盘价",
        "",
        "SlippageModel: 滑点0.1% + 佣金万三(最低5元) + 印花税千一(卖出)",
    ]

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    for i, item in enumerate(bt_items):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = item
        p.font.size = Pt(12 if item else 8)
        p.font.color.rgb = WHITE if item else RGBColor(0x44, 0x44, 0x44)
        p.font.name = "Consolas" if item.startswith("  ") else "Microsoft YaHei"
        p.space_after = Pt(1)

    # 右侧：BacktestResult 指标
    metrics = [
        ("核心绩效", SECONDARY, [
            "total_return — 累计收益率",
            "annual_return — 年化收益率",
            "sharpe_ratio — 夏普比率",
            "calmar_ratio — 卡玛比率",
            "max_drawdown — 最大回撤(%)",
        ]),
        ("交易统计", ACCENT, [
            "trade_count — 总交易次数",
            "win_rate — 胜率",
            "win_loss_ratio — 盈亏比",
            "profit_factor — 盈利因子",
            "avg_win / avg_loss — 均盈/均亏",
        ]),
    ]
    for i, (title, color, items) in enumerate(metrics):
        add_card(slide, 7.8, 1.8 + i * 2.5, 5.0, 2.3, title, items, title_color=color,
                 bg_color=RGBColor(0x25, 0x30, 0x50))

    add_footer(slide, "回测结果是历史模拟 · 不保证未来收益 · 过拟合/幸存者偏差可能导致结果失真")


# Slide 9: AI Agent
def add_agent_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "AI Agent 智能分析", "AI Agent — Market Diagnosis + Strategy Recommendation + Chat")

    # Agent 工作流
    steps = [
        ("📡\n数据加载", SECONDARY),
        ("🔍\n市场诊断", ACCENT),
        ("🎯\n策略推荐", ACCENT),
        ("🔧\n参数优化", WARNING),
        ("📋\n报告生成", WARNING),
        ("💬\n对话交互", RGBColor(0x9C, 0x27, 0xB0)),
    ]
    for i, (label, color) in enumerate(steps):
        left = 0.8 + i * 2.1
        shape = slide.shapes.add_shape(1, Inches(left), Inches(1.8), Inches(1.9), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(2, Inches(left + 1.9), Inches(2.15), Inches(0.2), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MID_GRAY
            arrow.line.fill.background()

    # 左侧：MarketRegimeDetector
    add_card(slide, 0.8, 3.2, 6.0, 2.0,
             "MarketRegimeDetector 市场状态诊断", [
                 "ADX判断趋势强度(>25有趋势, <20无趋势)",
                 "波动率分位分析(高/中/低三档)",
                 "RSI超买超卖检测",
                 "自动计算建议仓位比例",
                 "综合风险等级评估(低/中/高)",
                 "输出: 中文诊断描述 + 结构化数据",
             ], title_color=SECONDARY, bg_color=RGBColor(0x25, 0x30, 0x50))

    # 右侧：StrategySelector
    add_card(slide, 7.2, 3.2, 5.5, 2.0,
             "StrategySelector 策略自动推荐", [
                 "知识库: 5种策略 + 适用场景映射",
                 "趋势↑ → 双均线趋势跟踪",
                 "趋势↓ → 空仓/逆势反弹",
                 "震荡 → RSI均值回归",
                 "不明 → 多因子综合打分",
                 "波动率自适应微调参数",
             ], title_color=ACCENT, bg_color=RGBColor(0x25, 0x30, 0x50))

    # 对话 Agent
    add_card(slide, 0.8, 5.5, 11.7, 1.2,
             "ChatAgent 对话式交互", [
                 "规则引擎 + 关键词匹配实现NLU（无需LLM API，完全本地运行）",
                 "支持8种意图：分析/市场/推荐/对比/解释/优化/风控/回测",
                 "内置知识库：策略对比表 + 6个量化概念详解（夏普/回撤/MACD/RSI/凯利/卡玛）",
                 "对话历史管理 + 股票代码自动识别 + 快捷提问按钮",
             ], title_color=RGBColor(0x9C, 0x27, 0xB0), bg_color=RGBColor(0x25, 0x30, 0x50))

    add_footer(slide, "Agent基于统计规则，非LLM · 输出不构成交易建议")


# Slide 10: Visualization
def add_viz_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "可视化与参数优化", "Visualization & Optimization — 6 Charts + Grid Search + Overfitting")

    viz_items = [
        ("visualization/equity_curve.py", "权益曲线 vs 基准走势"),
        ("visualization/drawdown.py", "回撤水下曲线 + 最大回撤标注"),
        ("visualization/trades.py", "买卖点标记在K线图上"),
        ("visualization/pnl_histogram.py", "逐笔盈亏分布直方图"),
        ("visualization/signal_scatter.py", "信号强度 vs 未来收益散点"),
        ("visualization/dashboard.py", "四合一看板 (权益+回撤+信号+分布)"),
    ]
    for i, (file, desc) in enumerate(viz_items):
        top = 1.8 + i * 0.55
        add_textbox(slide, 0.8, top, 4.0, 0.4, file, font_size=11, color=SECONDARY, font_name="Consolas")
        add_textbox(slide, 5.0, top, 3.0, 0.4, "→ " + desc, font_size=12, color=LIGHT_GRAY)

    # 优化模块
    cards = [
        ("GridSearch 网格搜索", [
            "itertools.product 穷举组合",
            "自定义评分指标(夏普/收益/回撤)",
            "Top N 排序 + 最优参数输出",
            "支持静态参数锁定",
        ]),
        ("OverfittingDetector 过拟合检测", [
            "样本内(70%)网格搜索最优参数",
            "样本外(30%)验证绩效",
            "夏普衰减 > 50% → 过拟合警告",
            "Deflation Ratio(收缩比率)评估",
        ]),
        ("SensitivityAnalyzer 敏感性分析", [
            "OAT(One-at-a-time)逐个参数扫描",
            "Tornado 图数据生成",
            "识别高杠杆参数(敏感度排序)",
            "指导参数调优方向",
        ]),
    ]
    for i, (title, items) in enumerate(cards):
        add_card(slide, 0.8 + i * 4.1, 5.0, 3.8, 2.0, title, items,
                 title_color=WARNING, bg_color=RGBColor(0x25, 0x30, 0x50))

    add_footer(slide, "网格搜索加剧过拟合风险 · 须配合样本外验证 · 随机搜索/贝叶斯优化为备选方案")


# Slide 11: Project Structure
def add_structure_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "项目结构", "Project Structure — 60+ Python Files, 10 Modules")

    tree = """stock_quant_system/
    ├── agent/               🤖 AI量化Agent (市场诊断+策略推荐+对话)
    │   ├── quant_agent.py       主Agent控制器
    │   ├── market_regime.py     市场状态识别器
    │   ├── strategy_selector.py 策略自动推荐
    │   ├── auto_optimizer.py    参数自动优化
    │   ├── report_generator.py  报告生成器
    │   └── chat_agent.py       对话式Agent
    ├── config/              ⚙️ 全局配置 (DataConfig/BacktestConfig/RiskConfig)
    ├── data/                📡 数据层 (AKShare/YFinance/清洗/缓存)
    ├── factors/             📊 因子库 (12+指标 + 注册中心)
    ├── signals/             🎯 信号引擎 (趋势/反转/多因子/过滤链)
    ├── portfolio/           💰 仓位风控 (3模型 + 5规则 + 管理器)
    ├── backtest/            📈 回测引擎 (向量化模拟 + 指标 + 滑点)
    ├── optimization/        🔍 参数优化 (网格搜索 + 敏感性 + 过拟合)
    ├── visualization/       📉 可视化 (6种图表 + 综合看板)
    ├── live/                🔌 实盘接口抽象 (Order/Position/AccountInfo)
    ├── utils/               🛠️ 工具 (日志 + 异常体系)
    ├── examples/            📋 5个示例脚本
    ├── tests/               ✅ 测试
    └── app.py               🌐 Streamlit Web UI"""

    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.5))
    ctf = code_box.text_frame
    ctf.word_wrap = True
    for i, line in enumerate(tree.split("\n")):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.name = "Consolas"
        p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(0)

    # 右侧统计数据
    stats = [
        ("60+", "Python文件"),
        ("10", "功能模块"),
        ("12+", "技术因子"),
        ("5", "交易策略"),
        ("3", "仓位模型"),
        ("5", "风控规则"),
        ("6", "图表类型"),
        ("3", "Agent"),
    ]

    for i, (num, label) in enumerate(stats):
        row = i // 2
        col = i % 2
        left = 8.0 + col * 2.5
        top = 1.8 + row * 1.3

        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.2), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x25, 0x30, 0x50)
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.color.rgb = SECONDARY
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = "Microsoft YaHei"
        p2.alignment = PP_ALIGN.CENTER

    add_footer(slide, "模块化架构 · 每个模块可独立导入调用 · 松耦合设计")


# Slide 12: Tech Stack
def add_tech_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "技术栈", "Technology Stack — Python Open-Source Ecosystem")

    layers = [
        ("数据层", SECONDARY, [
            "AKShare — A股数据（免费无API Key）",
            "YFinance — 美股数据（备用）",
            "Pandas 3.x — 数据处理核心",
            "NumPy 2.x — 向量化数值计算",
        ]),
        ("计算层", ACCENT, [
            "Scipy — 统计分析（回归/检验）",
            "Statsmodels — 计量经济学模型",
            "Pandas-TA — 技术指标（备用）",
        ]),
        ("回测层", WARNING, [
            "VectorBT — 轻量向量化回测",
            "Backtrader — 事件驱动复杂回测",
            "自研 SimpleBacktestEngine",
        ]),
        ("界面层", RGBColor(0x9C, 0x27, 0xB0), [
            "Streamlit 1.58 — Web UI框架",
            "Matplotlib — 图表渲染",
            "python-pptx — PPT自动生成",
        ]),
    ]

    for i, (title, color, items) in enumerate(layers):
        add_card(slide, 0.8 + i * 3.2, 1.8, 2.9, 2.8, title, items,
                 title_color=color, bg_color=RGBColor(0x25, 0x30, 0x50))

    # 底部：依赖版本
    add_textbox(slide, 0.8, 5.2, 11, 0.4,
                "核心依赖: pandas>=2.2  numpy>=1.26  akshare>=1.14  matplotlib>=3.8  streamlit>=1.58",
                font_size=11, color=MID_GRAY)
    add_textbox(slide, 0.8, 5.6, 11, 0.4,
                "可选: vectorbt>=0.5  backtrader>=1.9  scipy>=1.12  scikit-learn",
                font_size=11, color=MID_GRAY)

    add_footer(slide, "全部使用Python开源工具链 · 无需商业授权")


# Slide 13: Risk & Limitations
def add_risk_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    add_slide_title(slide, "策略缺陷与过拟合风险", "Limitations & Overfitting — 不可忽视的风险")

    risks = [
        ("过拟合风险", DANGER, [
            "网格搜索在样本内表现好，样本外可能大幅衰减",
            "夏普衰减>50% 或样本外夏普<0 → 严重过拟合",
            "参数越多、搜索空间越大 → 过拟合越严重",
        ]),
        ("策略缺陷", WARNING, [
            "趋势策略: 震荡市频繁假信号 → 磨损本金",
            "均值回归: 强趋势市接飞刀 → 可能大亏",
            "多因子: 因子共线性 → 权重失效",
            "回测偏差: 前视偏差、幸存者偏差",
        ]),
        ("市场风险", DANGER, [
            "黑天鹅: 模型无法预测极端事件",
            "流动性: 小盘股滑点远超模型假设",
            "政策: A股T+1、涨跌停、停牌",
            "心理: 实盘亏损时的情绪干扰",
        ]),
        ("改进方向", ACCENT, [
            "样本外验证 + 交叉验证",
            "波动率过滤器（高波动暂停交易）",
            "策略组合降低单策略失效风险",
            "多周期参数平均（而非最优单值）",
            "Walk-Forward 滚动窗口分析",
        ]),
    ]

    for i, (title, color, items) in enumerate(risks):
        row = i // 2
        col = i % 2
        left = 0.8 + col * 6.2
        top = 1.8 + row * 2.7
        add_card(slide, left, top, 5.8, 2.4, title, items, title_color=color,
                 bg_color=RGBColor(0x25, 0x30, 0x50))

    add_footer(slide, "以上风险演示用于教学 · 实盘交易需专业风控和充分测试")


# Slide 14: Summary
def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)

    add_textbox(slide, 1.0, 0.8, 11, 0.6, "总结", font_size=36, bold=True, color=WHITE)
    line = slide.shapes.add_shape(1, Inches(1.0), Inches(1.5), Inches(11), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY
    line.line.fill.background()

    points = [
        ("✅", "完整的量化回测系统", "数据采集→因子计算→信号生成→仓位风控→回测→可视化"),
        ("✅", "AI Agent 智能分析", "市场诊断+策略推荐+参数优化+报告生成 全自动"),
        ("✅", "对话式交互", "自然语言查询分析，8种意图识别，无需LLM API"),
        ("✅", "模块化可扩展", "60+文件10大模块，因子注册中心支持热插拔"),
        ("✅", "Web UI + CLI", "Streamlit网页界面 (3功能标签) + 命令行脚本"),
        ("✅", "教学导向设计", "完整中文文档+注释+风险提示+过拟合检测"),
    ]

    for i, (icon, title, desc) in enumerate(points):
        top = 2.0 + i * 0.75
        add_textbox(slide, 1.0, top, 0.5, 0.5, icon, font_size=20, color=ACCENT)
        add_textbox(slide, 1.6, top, 3.5, 0.4, title, font_size=18, bold=True, color=WHITE)
        add_textbox(slide, 5.5, top, 7.0, 0.4, desc, font_size=14, color=LIGHT_GRAY)

    # 底部风险声明
    add_textbox(slide, 1.0, 6.3, 11, 0.6,
                "⚠️ 风险声明：本系统所有代码仅供教育及研究目的，不构成任何投资建议。"
                "过往业绩不代表未来表现。投资有风险，入市须谨慎。",
                font_size=13, color=DANGER, bold=True)

    add_textbox(slide, 1.0, 6.8, 11, 0.4,
                "2026 · Stock Quantitative Trading System · Built with Python", font_size=11, color=MID_GRAY)


if __name__ == "__main__":
    output = create_presentation()
    print(f"\nDone! PPT is ready at: {output}")
    print(f"Total slides: 14")
